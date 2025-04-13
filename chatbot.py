import streamlit as st
import google.generativeai as genai
from PIL import Image
import base64
import io
from docx import Document
import fitz  # PyMuPDF for PDF
from pptx import Presentation  # For PPTX files

# --- Setup ---
st.set_page_config(page_title="AI Chatbot", layout="centered")
genai.configure(api_key="")  # Replace with your Gemini API key
model = genai.GenerativeModel("models/gemini-1.5-flash")

# --- Session State ---
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "chat_session" not in st.session_state:
    st.session_state.chat_session = model.start_chat(history=[])

# --- Link to the external CSS file ---
with open("styles.css") as f:
    st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

# --- Title ---
st.markdown("<h1 style='text-align:center;'>💬 AI Chatbot</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align:center; color: gray;'>Let's Create Something Amazing Together!</p>", unsafe_allow_html=True)

# --- Quick Suggestions ---
suggestions = [
    "Help me write an e-mail ✉️", "Translate a sentence 🌍",
    "How to build a resume 📋 ", "Summarize a concept 🗞️", "Give me coding help 👨‍💻"
]
st.markdown("#### 💡 Quick Suggestions to Start with:")
cols = st.columns(len(suggestions))
for i, s in enumerate(suggestions):
    if cols[i].button(s):
        st.session_state.chat_history.append({"role": "user", "text": s})
        response = st.session_state.chat_session.send_message(s)
        st.session_state.chat_history.append({"role": "assistant", "text": response.text})
        st.rerun()

# --- Display Chat History ---
for msg in st.session_state.chat_history:
    with st.chat_message(msg["role"]):
        st.markdown(msg["text"])
        if msg.get("image"):
            st.image(msg["image"], width=300)

# --- Chat Input ---
st.markdown('<div class="chat-input-container">', unsafe_allow_html=True)
with st.form(key="chat_form", clear_on_submit=True):
    uploaded_file = st.file_uploader(
        label="Browse file (image/pdf/docx/pptx)",
        type=["jpg", "jpeg", "png", "pdf", "docx", "pptx"],
        label_visibility="visible"
    )

    input_cols = st.columns([0.85, 0.15])
    with input_cols[0]:
        prompt = st.text_area(
            "Ask something...",
            height=80,
            label_visibility="collapsed",
            placeholder="Type your message here...",
            key="text_area_input"
        )
    with input_cols[1]:
        submit_btn = st.form_submit_button("➤")
st.markdown("</div>", unsafe_allow_html=True)

# --- Handle Chat Submission ---
if submit_btn and (prompt or uploaded_file):
    user_msg = {"role": "user", "text": prompt}
    user_parts = []

    if uploaded_file:
        file_bytes = uploaded_file.read()
        mime_type = uploaded_file.type
        file_name = uploaded_file.name

        if mime_type.startswith("image/"):
            user_msg["image"] = file_bytes
            user_parts.append({
                "inline_data": {
                    "mime_type": mime_type,
                    "data": base64.b64encode(file_bytes).decode()
                }
            })

        elif mime_type == "application/pdf":
            extracted_text = ""
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf_doc:
                for page in pdf_doc:
                    extracted_text += page.get_text()
            user_parts.append({"text": f"Here's the extracted text from the PDF:\n{extracted_text.strip()}"})

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            uploaded_file.seek(0)
            docx_bytes = io.BytesIO(uploaded_file.read())
            doc = Document(docx_bytes)
            full_text = "\n".join([para.text for para in doc.paragraphs])
            user_parts.append({"text": f"Here's the extracted text from the DOCX file:\n{full_text.strip()}"})

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            try:
                ppt = Presentation(io.BytesIO(file_bytes))
                ppt_text = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            ppt_text += f"Slide {ppt.slides.index(slide)+1}: {shape.text}\n"
                user_parts.append({"text": f"Here's the extracted text from the PPTX file:\n{ppt_text.strip()}"})
            except Exception as e:
                st.error(f"An error occurred while processing the PowerPoint file: {str(e)}")

    # Add prompt if available
    if prompt:
        user_parts.append({"text": prompt})

    # Save message to chat history
    st.session_state.chat_history.append(user_msg)

    # Get model response
    response = st.session_state.chat_session.send_message(user_parts)
    st.session_state.chat_history.append({"role": "assistant", "text": response.text})
    st.rerun()
